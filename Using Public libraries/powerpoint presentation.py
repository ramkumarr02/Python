# -*- coding: utf-8 -*-
"""
Created on Sun Jan 08 16:08:22 2017

@author: ramrajagopalan
"""

from pptx import presentation

deck = Presentation()

title_slide_layout = deck.slide_layouts[0]

slide = deck.slides.add_slide(title_slide_layout)

title = slide.shapes.title

title.text = "Hi there"

deck.save("test2.pptx")

